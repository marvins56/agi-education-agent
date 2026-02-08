"""PPTX presentation parser using python-pptx.

Extracts text from slides including titles, body text boxes,
table cells, and speaker notes.
"""

import logging

from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger(__name__)


def _table_to_markdown(table) -> str:
    """Convert a python-pptx table to a markdown-formatted table string."""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cell_text = cell.text.strip().replace("|", "\\|")
            cells.append(cell_text)
        rows.append(cells)

    if not rows:
        return ""

    lines = []
    lines.append("| " + " | ".join(rows[0]) + " |")
    lines.append("| " + " | ".join("---" for _ in rows[0]) + " |")
    for row in rows[1:]:
        while len(row) < len(rows[0]):
            row.append("")
        lines.append("| " + " | ".join(row[: len(rows[0])]) + " |")

    return "\n".join(lines)


class PptxParser:
    """Extract text and metadata from PPTX presentation files."""

    def parse(self, file_path: str) -> str:
        """Extract all text from a PPTX file.

        Returns text organized by slide with titles, body text,
        tables as markdown, and speaker notes.
        """
        try:
            prs = Presentation(file_path)
        except Exception as exc:
            logger.error("Failed to open PPTX file %s: %s", file_path, exc)
            return ""

        parts = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            try:
                slide_parts = self._extract_slide(slide, slide_num)
                if slide_parts:
                    parts.append("\n".join(slide_parts))
            except Exception as exc:
                logger.debug("Error processing slide %d: %s", slide_num, exc)
                continue

        return "\n\n".join(parts)

    def parse_with_metadata(self, file_path: str) -> dict:
        """Extract text and metadata from a PPTX file."""
        try:
            prs = Presentation(file_path)
        except Exception as exc:
            logger.error("Failed to open PPTX file %s: %s", file_path, exc)
            return {"text": "", "pages": 0, "metadata": {}}

        parts = []
        slide_count = 0
        has_notes = False
        table_count = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_count += 1
            try:
                slide_parts = self._extract_slide(slide, slide_num)
                if slide_parts:
                    parts.append("\n".join(slide_parts))

                # Count tables in this slide
                for shape in slide.shapes:
                    try:
                        if shape.has_table:
                            table_count += 1
                    except Exception:
                        continue

                # Check for notes
                try:
                    if (
                        slide.has_notes_slide
                        and slide.notes_slide
                        and slide.notes_slide.notes_text_frame
                    ):
                        notes_text = slide.notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            has_notes = True
                except Exception:
                    pass

            except Exception as exc:
                logger.debug("Error processing slide %d: %s", slide_num, exc)
                continue

        text = "\n\n".join(parts)

        return {
            "text": text,
            "pages": slide_count,
            "metadata": {
                "slide_count": slide_count,
                "has_notes": has_notes,
                "table_count": table_count,
            },
        }

    def _extract_slide(self, slide, slide_num: int) -> list[str]:
        """Extract all content from a single slide."""
        parts: list[str] = []

        # Get slide title
        title = self._get_slide_title(slide)
        title_display = title if title else "Untitled"
        parts.append(f"## Slide {slide_num}: {title_display}")

        # Extract text from all shapes
        body_texts = []
        tables = []

        for shape in slide.shapes:
            try:
                # Skip the title shape since we already captured it
                if shape.has_text_frame:
                    if shape == getattr(slide.shapes, "title", None):
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            body_texts.append(text)

                if shape.has_table:
                    md_table = _table_to_markdown(shape.table)
                    if md_table:
                        tables.append(md_table)
            except Exception as exc:
                logger.debug(
                    "Error processing shape on slide %d: %s", slide_num, exc
                )
                continue

        if body_texts:
            parts.extend(body_texts)

        if tables:
            parts.extend(tables)

        # Extract speaker notes
        try:
            if (
                slide.has_notes_slide
                and slide.notes_slide
                and slide.notes_slide.notes_text_frame
            ):
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    parts.append(f"[Notes: {notes_text}]")
        except Exception as exc:
            logger.debug("Error reading notes on slide %d: %s", slide_num, exc)

        return parts

    def _get_slide_title(self, slide) -> str:
        """Extract the title text from a slide, if present."""
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                return slide.shapes.title.text_frame.text.strip()
        except Exception:
            pass

        # Fallback: look for a shape with a title placeholder type
        try:
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 0:  # Title placeholder
                    return shape.text.strip()
        except Exception:
            pass

        return ""
